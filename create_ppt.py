from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 62)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 175, 55)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(136, 136, 136)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 62)
    background.line.fill.background()
    
    # 标题栏
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(15, 15, 46)
    title_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 175, 55)
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(224, 224, 224)
        p.space_before = Pt(12)
        p.line_spacing = 1.5
    
    return slide

def add_two_col_slide(prs, title, left_title, left_content, right_title, right_content):
    """添加两列布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 62)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 175, 55)
    
    # 左列标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(243, 156, 18)
    
    # 左列内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(224, 224, 224)
        p.space_before = Pt(8)
    
    # 右列标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 152, 219)
    
    # 右列内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(224, 224, 224)
        p.space_before = Pt(8)
    
    return slide

# ===== 开始创建PPT =====

# 第1页：封面
add_title_slide(prs, "🎯 计策系统设计方案", "第三阶段：基于338位英雄数据的SRPG核心系统建议\n\n数据来源：梦幻模拟战208位 + 天地劫150+位 + 铃兰之剑82位 + 三国志战棋版126位\n\n2026年3月14日")

# 第2页：目录
add_content_slide(prs, "📋 目录", [
    "1. 核心洞察：计策的本质是\"战中成长\"",
    "2. 设计框架：概念 → 结构 → 细节 → 体验",
    "3. 三大计策类型：时机型 / 环境型 / 数值型",
    "4. 成长维度：数值累加 / 优势造成 / 适应型",
    "5. 限制体系：CD / 能量 / 代价 / 条件",
    "6. 关卡×计策：场景化设计对应表",
    "7. 关键原则：设计的红线与底线"
])

# 第3页：核心洞察
add_content_slide(prs, "💡 核心洞察：计策的本质", [
    "\"战中成长\"——让玩家在战斗过程中持续变强",
    "",
    "不同于被动技能，计策的本质是「让玩家在战斗过程中持续成长、积累优势的能力」。",
    "这是与传统SRPG最本质的差异——不是静态的数值比拼，而是动态的优势建构。",
    "",
    "三种成长模式：",
    "• 数值累加型：每次行动叠加层数，后期爆发（例：司马懿\"鹰视狼顾\"）",
    "• 优势造成型：通过计策创造地形/时机优势，滚雪球扩大领先",
    "• 适应型成长：根据战场变化调整策略，越打越顺的掌控感"
])

# 第4页：设计框架
add_content_slide(prs, "🏗️ 设计层级框架", [
    "概念层 (Why)：计策系统的存在意义",
    "  → 解决\"开局定胜负\"的枯燥感 / 提供战中成长与变数 / 让玩家掌握\"改写规则\"的权力",
    "",
    "结构层 (What)：计策系统的骨架",
    "  → 三大类计策：时机/环境/数值 / 成长维度：累加/优势/适应 / 限制体系：CD/能量/代价/条件",
    "",
    "细节层 (How)：具体数值与规则",
    "  → 数值边界：单体≤1.8倍，AOE≤0.7倍 / 持续时间：1-3回合最佳 / 触发条件设计",
    "",
    "体验层 (Feel)：玩家感受到什么",
    "  → \"我在变强\"的正反馈 / 每回合都有新选择 / 逆境翻盘的爽快感 / 策略深度的掌控感"
])

# 第5页：时机型计策
add_content_slide(prs, "⚡ 类型一：时机型计策", [
    "利用战场时机创造优势的计策，考验玩家对节奏的把握。",
    "",
    "【先发制人】战斗开始时触发 → 首回合移动力+2，先攻",
    "   适合快攻流派，抢占先机",
    "",
    "【乘胜追击】击杀敌人后触发 → 再行动一次，伤害+20%",
    "   滚雪球核心，连续收割",
    "",
    "【背水一战】生命值<30%触发 → 攻击+50%，受到伤害+30%",
    "   高风险高回报，绝境翻盘",
    "",
    "【以逸待劳】待机后触发 → 下次攻击伤害+40%，必中",
    "   鼓励策略性待机",
    "",
    "关键洞察：触发条件的平衡——太容易触发会滥用，太难则无人问津"
])

# 第6页：环境型计策
add_two_col_slide(prs, "🌍 类型二：环境型计策",
    "🔥 地形改变",
    [
        "【火烧连营】",
        "区域持续灼烧，",
        "敌人经过受伤",
        "",
        "【水淹七军】",
        "制造水域，",
        "减速+易伤",
        "",
        "【八门金锁】",
        "生成阻挡地形，",
        "改变路径"
    ],
    "⚡ 领域效果",
    [
        "【八卦阵】",
        "范围内友方",
        "闪避+30%",
        "",
        "【奇门遁甲】",
        "敌人进入时",
        "随机debuff",
        "",
        "【桃园结义】",
        "范围内友方",
        "每回合回血"
    ]
)

# 第7页：数值型计策
add_content_slide(prs, "📊 类型三：数值型计策", [
    "直接改变数值关系，简单粗暴但效果显著。",
    "",
    "【破釜沉舟】爆发型：消耗50%当前生命，伤害×2.5",
    "   适用场景：Boss斩杀/绝境反击",
    "",
    "【稳扎稳打】成长型：每回合攻击+10%，可叠加5层",
    "   适用场景：持久战/后期核心",
    "",
    "【以弱胜强】克制型：对高等级敌人伤害+50%",
    "   适用场景：越级挑战",
    "",
    "【援护】防御型：替相邻友方承受伤害，减伤30%",
    "   适用场景：保护核心输出",
    "",
    "⚠️ 数值红线：单体计策倍率不超过1.8倍，AOE不超过0.7倍"
])

# 第8页：成长维度
add_two_col_slide(prs, "📈 成长维度详解",
    "📊 数值累加型",
    [
        "每次行动叠加层数，",
        "后期爆发",
        "",
        "• 线性成长：每回合+10%",
        "• 指数成长：层数平方增长",
        "• 阶梯成长：每3回合质变",
        "",
        "代表：司马懿鹰视狼顾、",
        "诸葛亮五层被动"
    ],
    "🎯 优势造成型",
    [
        "创造环境优势，",
        "滚雪球扩大",
        "",
        "• 地形优势：火海/水域/高地",
        "• 位置优势：包夹/夹击/围攻",
        "• 时机优势：先手/再动/打断",
        "",
        "代表：周瑜火烧连营、",
        "曹操乱世奸雄"
    ]
)

# 第9页：限制体系
add_content_slide(prs, "⚖️ 限制体系设计", [
    "计策必须有合理的限制，否则会变成滥用的\"万能解\"。",
    "",
    "⏱️ CD限制：使用后进入冷却",
    "   短CD(1-2回合) / 中CD(3-4回合) / 长CD(5+回合)",
    "",
    "⚡ 能量限制：消耗能量释放",
    "   小计策20点 / 中计策50点 / 大计策100点",
    "",
    "💀 代价限制：释放有代价",
    "   生命值消耗 / 资源消耗 / 负面状态",
    "",
    "🎯 条件限制：特定条件触发",
    "   站位要求 / 状态要求 / 时机要求",
    "",
    "限制设计黄金比例：效果强度 ∝ 限制强度"
])

# 第10页：关卡对应
add_content_slide(prs, "🎮 关卡×计策对应设计", [
    "不同关卡场景应该推荐不同的计策策略，形成教学-应用闭环。",
    "",
    "BOSS战 → 数值型爆发 + 时机控制",
    "   挑战：高血量/高伤害/阶段转换 → 学习资源管理与时机把握",
    "",
    "防守关 → 环境控制 + 区域增益",
    "   挑战：保护目标/多波次敌人 → 学习地形利用与布局",
    "",
    "限时关 → 先攻/再动 + 群体增益",
    "   挑战：回合限制/速战速决 → 学习效率最大化",
    "",
    "突围关 → 位移/机动 + 控制",
    "   挑战：从起点到达终点 → 学习机动与控场",
    "",
    "群殴关 → AOE伤害 + 连锁反应",
    "   挑战：大量杂兵 → 学习群体控制"
])

# 第11页：关键原则
add_two_col_slide(prs, "⚠️ 关键设计原则",
    "✅ 应该做",
    [
        "• 让玩家每回合都有选择",
        "  保持决策频率和参与感",
        "",
        "• 提供多种解题思路",
        "  鼓励创造性策略",
        "",
        "• 建立清晰的反馈回路",
        "  让玩家知道计策生效",
        "",
        "• 平衡风险与收益",
        "  强效果配强限制"
    ],
    "❌ 不要做",
    [
        "• 设计\"必带\"计策",
        "  破坏多样性",
        "",
        "• 让计策替代基础战斗",
        "  计策是补充不是替代",
        "",
        "• 设计过于复杂的效果",
        "  玩家记不住就不会用",
        "",
        "• 忽视PVE与PVP的平衡",
        "  两种模式需要不同设计"
    ]
)

# 第12页：数值红线
add_content_slide(prs, "🚨 数值红线（基于338位英雄数据）", [
    "基于对4款主流SRPG（338位英雄/武将）的深度分析，",
    "以下数值边界是保持战斗平衡的红线：",
    "",
    "• 单体计策倍率 ≤ 1.8倍基础伤害",
    "  超过此值将破坏基础战斗平衡",
    "",
    "• AOE计策倍率 ≤ 0.7倍基础伤害",
    "  群体伤害需要更高的代价或限制",
    "",
    "• 再动机制：每回合最多1次或击杀触发",
    "  无限再动将彻底破坏回合制节奏",
    "",
    "• 控制时长：眩晕/冰冻 ≤ 2回合",
    "  过长的控制等于罚站，体验极差",
    "",
    "• 增益上限：单项属性提升 ≤ 50%",
    "  过高的增益会让其他系统失效"
])

# 第13页：总结
add_content_slide(prs, "🎯 核心结论", [
    "",
    "",
    "计策系统的设计核心是「战中成长」",
    "",
    "通过时机、环境、数值三个维度",
    "",
    "让玩家在战斗中持续积累优势",
    "",
    "",
    "让每一回合都充满期待",
    "让每一场战斗都独一无二",
    "",
    "",
    "—",
    "基于梦幻模拟战208位 + 天地劫150+位 + 铃兰之剑82位 + 三国志战棋版126位英雄数据"
])

# 保存文件
output_path = "/root/.openclaw/workspace/portfolio-blog/research/srpg-analysis/phase3-strategy-design.pptx"
prs.save(output_path)
print(f"PPTX文件已保存至: {output_path}")
print(f"共创建 {len(prs.slides)} 页幻灯片")
